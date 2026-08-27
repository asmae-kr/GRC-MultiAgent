"""
generate_report_pptx.py
Genere le rapport PowerPoint (6 slides) d'evaluation fournisseur GRC directement
a partir du JSON de resultat produit par les agents - sans template a remplir.

Structure (identique a GRC_Template_6Slides.pptx) :
  1. Titre
  2. Resume executif
  3. Vue d'ensemble de la conformite (donut + KPI)
  4. Conforme vs Non conforme (deux cartes vert/rouge)
  5. Principaux risques identifies
  6. Recommandations priorisees

Usage (depuis watcher.py), juste apres generate_excel_from_template(...) :

    from generate_report_pptx import generate_report_pptx
    generate_report_pptx(result, output_path="Rapports/Rapport_Acme.pptx",
                          supplier_name="Acme Corp", evaluation_date="27/08/2026")

`result` est le dict JSON sauvegarde par write_result(), de forme :
{ "report": { "executive_summary": ..., "total_questions": ..., "conforme_count": ...,
  "partiellement_conforme_count": ..., "non_conforme_count": ...,
  "compliance_results": [ {requirement_id, criterion, status, score, justification, importance}, ... ] } }
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

ICONS_DIR = Path(__file__).parent / "icons"

# ---------------------------------------------------------------- palette
NAVY = RGBColor(0x1E, 0x27, 0x61)
NAVY_DARK = RGBColor(0x14, 0x1B, 0x4D)
ICE = RGBColor(0xCA, 0xDC, 0xFC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
OFFWHITE = RGBColor(0xF4, 0xF6, 0xFB)
GRAY = RGBColor(0x5B, 0x64, 0x79)
RED = RGBColor(0xC0, 0x39, 0x2B)
RED_BG = RGBColor(0xFB, 0xE7, 0xE4)
AMBER = RGBColor(0xD9, 0x8E, 0x04)
GREEN = RGBColor(0x1E, 0x8A, 0x5F)
GREEN_BG = RGBColor(0xE4, 0xF3, 0xEC)
LIGHTLINE = RGBColor(0xE3, 0xE7, 0xF1)

EMU_IN = 914400
SLIDE_W = 13.333
SLIDE_H = 7.5


def _in(v):
    return Inches(v)


def _add_rect(slide, x, y, w, h, color, radius=None, line_color=None, line_w=None):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, _in(x), _in(y), _in(w), _in(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    if line_color:
        shp.line.color.rgb = line_color
        shp.line.width = Pt(line_w or 1)
    else:
        shp.line.fill.background()
    shp.shadow.inherit = False
    if radius:
        try:
            shp.adjustments[0] = radius
        except Exception:
            pass
    return shp


def _add_ellipse(slide, x, y, w, h, color):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, _in(x), _in(y), _in(w), _in(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def _add_line(slide, x, y, w, color, weight=1.0):
    ln = slide.shapes.add_connector(1, _in(x), _in(y), _in(x + w), _in(y))
    ln.line.color.rgb = color
    ln.line.width = Pt(weight)
    return ln


def _add_icon(slide, name, x, y, w, h):
    path = ICONS_DIR / f"{name}.png"
    if path.exists():
        slide.shapes.add_picture(str(path), _in(x), _in(y), _in(w), _in(h))


def _add_text(slide, text, x, y, w, h, size=12, color=NAVY_DARK, bold=False,
              italic=False, font="Calibri", align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
              line_spacing=1.0, wrap=True):
    box = slide.shapes.add_textbox(_in(x), _in(y), _in(w), _in(h))
    tf = box.text_frame
    tf.word_wrap = wrap
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    if line_spacing != 1.0:
        p.line_spacing = line_spacing
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.name = font
    r.font.color.rgb = color
    return box


def _add_rich_text(slide, runs, x, y, w, h, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE, size=12):
    box = slide.shapes.add_textbox(_in(x), _in(y), _in(w), _in(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    for text, opts in runs:
        r = p.add_run()
        r.text = text
        r.font.size = Pt(opts.get("size", size))
        r.font.bold = opts.get("bold", False)
        r.font.name = "Calibri"
        r.font.color.rgb = opts.get("color", NAVY_DARK)
    return box


def _blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def _footer(slide, page_num, label):
    _add_text(slide, label, 0.5, SLIDE_H - 0.42, 6, 0.3, size=9, color=GRAY)
    _add_text(slide, str(page_num), SLIDE_W - 1.0, SLIDE_H - 0.42, 0.5, 0.3, size=9, color=GRAY, align=PP_ALIGN.RIGHT)


def _section_tag(slide, text):
    box = _add_text(slide, text.upper(), 0.6, 0.45, 8, 0.35, size=12, bold=True, color=NAVY)
    r = box.text_frame.paragraphs[0].runs[0]
    r._r.get_or_add_rPr().set('spc', '200')


def _slide_title(slide, text):
    _add_text(slide, text, 0.6, 0.78, 11.5, 0.7, size=30, bold=True, color=NAVY_DARK, font="Cambria")


def _fmt_pct(p):
    return f"{int(p)}%" if p == int(p) else f"{round(p, 1)}%"


def _truncate(text, max_chars):
    text = text or ""
    if len(text) <= max_chars:
        return text
    return text[:max_chars - 1].rstrip() + "\u2026"


# ============================================================================
def generate_report_pptx(result: dict, output_path: str, supplier_name: str = "[Nom du fournisseur]",
                          evaluation_date: str = None, prepared_by: str = "Equipe GRC"):

    from datetime import datetime
    if evaluation_date is None:
        evaluation_date = datetime.now().strftime("%d/%m/%Y")

    report = result.get("report", result)
    total = report.get("total_questions", 0) or 0
    conforme = report.get("conforme_count", 0) or 0
    partiel = report.get("partiellement_conforme_count", 0) or 0
    non_conforme = report.get("non_conforme_count", 0) or 0
    exec_summary = report.get("executive_summary", "") or "[Resume non disponible]"
    results = report.get("compliance_results", [])

    pct_conf = round((conforme / total) * 100, 1) if total else 0
    pct_non = round((non_conforme / total) * 100, 1) if total else 0

    conformes_items = [r.get("criterion", "") for r in results if r.get("status") == "Conforme"]
    non_conformes_items = [(r.get("requirement_id", ""), r.get("criterion", ""))
                            for r in results if r.get("status") == "Non conforme"]

    def _crit_level(r):
        """Normalise le niveau de criticite quel que soit le nom du champ
        (importance / criticality / Criticality) ou l'echelle utilisee
        (Major/Standard/Minor, High/Medium/Low, etc.)."""
        raw = r.get("importance") or r.get("criticality") or r.get("Criticality") or ""
        val = str(raw).strip().lower()
        if val in ("major", "high", "critical", "critique", "eleve", "\u00e9lev\u00e9"):
            return "Major"
        if val in ("minor", "low", "faible", "mineur"):
            return "Minor"
        return "Standard"  # standard / medium / moyen / valeur inconnue -> defaut

    majors_nc = [r for r in results if r.get("status") == "Non conforme" and _crit_level(r) == "Major"]
    standards_nc = [r for r in results if r.get("status") == "Non conforme" and _crit_level(r) == "Standard"]
    minors_nc = [r for r in results if r.get("status") == "Non conforme" and _crit_level(r) == "Minor"]

    prs = Presentation()
    prs.slide_width = Emu(int(SLIDE_W * EMU_IN))
    prs.slide_height = Emu(int(SLIDE_H * EMU_IN))
    footer_label = f"Rapport d'evaluation \u2014 {supplier_name}"
    page = [1]  # compteur mutable (liste pour etre modifie dans les closures)

    def _foot(slide):
        _footer(slide, page[0], footer_label)
        page[0] += 1

    # ---------------------------------------------------------- SLIDE 1: TITRE
    s = _blank_slide(prs)
    _bg(s, NAVY_DARK)
    _add_ellipse(s, 9.6, 3.0, 6, 6, NAVY)
    _add_ellipse(s, 10.6, 4.0, 4.2, 4.2, NAVY_DARK)
    _add_icon(s, "shield-white", 0.65, 0.7, 0.5, 0.5)
    _add_text(s, "GRC  \u00b7  VENDOR RISK ASSESSMENT", 1.25, 0.68, 8, 0.5, size=13, color=ICE, bold=True)
    _add_text(s, "Rapport d'evaluation\nde risque fournisseur", 0.65, 2.55, 10.5, 2.3,
              size=44, color=WHITE, bold=True, font="Cambria", line_spacing=1.08)
    _add_text(s, "Synthese de conformite, cartographie des risques et plan d'action",
              0.67, 4.55, 9.5, 0.5, size=16, color=ICE, italic=True)
    _add_line(s, 0.67, 5.35, 2.2, ICE, 1.5)
    _add_rich_text(s, [("Fournisseur Evalue :  ", dict(bold=True, color=ICE)), (supplier_name, dict(color=WHITE))],
                    0.67, 5.6, 8, 0.35, size=13)
    _add_rich_text(s, [("Date de l'evaluation :  ", dict(bold=True, color=ICE)), (evaluation_date, dict(color=WHITE))],
                    0.67, 5.98, 8, 0.35, size=13)
    _add_rich_text(s, [("Prepare par :  ", dict(bold=True, color=ICE)), (prepared_by, dict(color=WHITE))],
                    0.67, 6.36, 8, 0.35, size=13)
    _add_text(s, "CONFIDENTIEL \u2014 USAGE INTERNE", 0.67, SLIDE_H - 0.6, 6, 0.3, size=9, color=GRAY)

    # ---------------------------------------------------- SLIDE 2: RESUME EXECUTIF
    s = _blank_slide(prs)
    _section_tag(s, "Resume executif")
    _slide_title(s, "Niveau de risque global")

    _add_rect(s, 0.6, 1.85, 5.0, 4.85, OFFWHITE, radius=0.03)
    _add_icon(s, "file-text-navy", 0.95, 2.15, 0.4, 0.4)
    _add_text(s, "Resume", 1.5, 2.15, 3, 0.4, size=14, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    _add_text(s, exec_summary, 0.95, 2.75, 4.3, 3.7, size=13, color=GRAY, line_spacing=1.35)

    stat_color = GREEN if pct_conf >= 70 else (AMBER if pct_conf >= 40 else RED)
    stats = [
        (str(total), "Exigences evaluees", GRAY, "clipboard-navy"),
        (_fmt_pct(pct_conf), "Conformite globale", stat_color, "x-circle-red"),
        (str(conforme), "Conformes", GREEN, "check-circle-navy"),
        (str(non_conforme), "Non conformes", RED, "alert-triangle-navy"),
    ]
    gx, gy, cw, ch, gap = 5.95, 1.85, 3.15, 2.35, 0.2
    for i, (val, label, color, icon) in enumerate(stats):
        col, row = i % 2, i // 2
        x = gx + col * (cw + gap)
        y = gy + row * (ch + gap)
        _add_rect(s, x, y, cw, ch, WHITE, radius=0.04, line_color=LIGHTLINE, line_w=0.75)
        _add_icon(s, icon, x + 0.25, y + 0.25, 0.38, 0.38)
        _add_text(s, val, x + 0.15, y + 0.65, cw - 0.3, 1.0, size=42, bold=True, color=color, font="Cambria")
        _add_text(s, label, x + 0.15, y + ch - 0.55, cw - 0.3, 0.4, size=12, color=GRAY)

    _foot(s)

    # ------------------------------------------------- SLIDE 3: VUE D'ENSEMBLE
    s = _blank_slide(prs)
    _section_tag(s, "Resultats de conformite")
    _slide_title(s, "Vue d'ensemble de la conformite")

    chart_data = CategoryChartData()
    chart_data.categories = ["Conforme", "Partiellement conforme", "Non conforme"]
    vals = (conforme, partiel, non_conforme)
    chart_data.add_series("Statut", vals if sum(vals) > 0 else (1, 1, 1))
    gframe = s.shapes.add_chart(XL_CHART_TYPE.DOUGHNUT, _in(1.4), _in(1.95), _in(5.6), _in(4.9), chart_data)
    chart = gframe.chart
    chart.has_title = False
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    plot = chart.plots[0]
    plot.has_data_labels = True
    plot.data_labels.show_percentage = True
    plot.data_labels.show_value = False
    for i, point in enumerate(plot.series[0].points):
        point.format.fill.solid()
        point.format.fill.fore_color.rgb = [GREEN, AMBER, RED][i]

    pct_partiel = round((partiel / total) * 100, 1) if total else 0

    kpis = [
        (_fmt_pct(pct_conf), "Conforme", GREEN, str(conforme)),
        (_fmt_pct(pct_partiel), "Partiellement conforme", AMBER, str(partiel)),
        (_fmt_pct(pct_non), "Non conforme", RED, str(non_conforme)),
    ]
    ky = 1.95
    for pct_val, label, color, count_val in kpis:
        _add_rect(s, 7.6, ky, 5.1, 1.45, OFFWHITE, radius=0.08)
        _add_text(s, pct_val, 7.9, ky + 0.15, 2.1, 1.15, size=36, bold=True, color=color, font="Cambria", anchor=MSO_ANCHOR.MIDDLE)
        _add_text(s, label, 10.0, ky + 0.25, 2.5, 0.5, size=14, bold=True, color=NAVY_DARK)
        _add_text(s, f"{count_val} / {total} exigences", 10.0, ky + 0.75, 2.5, 0.4, size=11.5, color=GRAY)
        ky += 1.65

    _foot(s)

    # ------------------------------------------- SLIDE 4+: CONFORME VS NON CONFORME (pagine)
    colW, colGap, colY, colH, startX = 5.75, 0.4, 1.9, 4.75, 0.6

    # espace dispo sous l'entete de carte pour la liste = colH - 1.55 (entete) - 0.15 (marge basse)
    LIST_TOP_OFFSET = 1.55
    AVAILABLE_H = colH - LIST_TOP_OFFSET - 0.15
    ROW_H = 0.36                      # ligne compacte pour lister un maximum d'items
    PER_PAGE = max(1, int(AVAILABLE_H // ROW_H))   # nb d'items par colonne et par page

    def _paginate(items, per_page):
        if not items:
            return [[]]
        return [items[i:i + per_page] for i in range(0, len(items), per_page)]

    conforme_pages = _paginate(conformes_items, PER_PAGE)
    non_conforme_pages = _paginate(non_conformes_items, PER_PAGE)
    num_pages = max(len(conforme_pages), len(non_conforme_pages))

    for page_idx in range(num_pages):
        s = _blank_slide(prs)
        _section_tag(s, "Resultats de conformite")
        title = "Conforme vs Non conforme"
        if num_pages > 1:
            title += f"  ({page_idx + 1}/{num_pages})"
        _slide_title(s, title)

        conforme_chunk = conforme_pages[page_idx] if page_idx < len(conforme_pages) else []
        non_conforme_chunk = non_conforme_pages[page_idx] if page_idx < len(non_conforme_pages) else []

        def _panel(x, title_txt, count_txt, pct_txt, icon, border, bg, chunk, is_conforme, items_all):
            _add_rect(s, x, colY, colW, colH, WHITE, radius=0.03, line_color=border, line_w=1.5)
            _add_ellipse(s, x + 0.35, colY + 0.35, 0.7, 0.7, border)
            _add_icon(s, icon, x + 0.5, colY + 0.5, 0.4, 0.4)
            _add_text(s, title_txt, x + 1.2, colY + 0.3, colW - 2.2, 0.45, size=18, bold=True)
            _add_text(s, count_txt, x + 1.2, colY + 0.72, colW - 2.2, 0.35, size=12, color=GRAY)
            _add_text(s, pct_txt, x + colW - 2.0, colY + 0.25, 1.65, 0.85, size=28, bold=True,
                       color=border, font="Cambria", align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE, wrap=False)
            _add_line(s, x + 0.35, colY + 1.3, colW - 0.7, LIGHTLINE, 1)

            if not items_all:
                _add_text(s, "Aucune exigence dans cette categorie.", x + 0.35, colY + 1.7, colW - 0.7, 0.4,
                           size=12, italic=True, color=GRAY)
                return

            cy = colY + LIST_TOP_OFFSET
            for it in chunk:
                _add_ellipse(s, x + 0.32, cy + 0.02, 0.24, 0.24, bg)
                _add_icon(s, "check-circle-green" if is_conforme else "x-circle-red", x + 0.36, cy + 0.06, 0.16, 0.16)
                if is_conforme:
                    _add_text(s, _truncate(it, 60), x + 0.68, cy, colW - 1.0, ROW_H, size=10.5,
                               anchor=MSO_ANCHOR.MIDDLE, wrap=False)
                else:
                    rid, crit = it
                    _add_rich_text(s, [(rid + "  ", dict(bold=True, color=GRAY, size=9.5)),
                                        (_truncate(crit, 52), dict(color=NAVY_DARK))],
                                    x + 0.68, cy, colW - 1.0, ROW_H, size=10.5)
                cy += ROW_H

        _panel(startX, "Conforme", f"{conforme} / {total} exigences", _fmt_pct(pct_conf),
               "check-white", GREEN, GREEN_BG, conforme_chunk, True, conformes_items)
        _panel(startX + colW + colGap, "Non conforme", f"{non_conforme} / {total} exigences", _fmt_pct(pct_non),
               "x-white", RED, RED_BG, non_conforme_chunk, False, non_conformes_items)

        _foot(s)

    # ------------------------------------------------- SLIDE 5: PRINCIPAUX RISQUES
    s = _blank_slide(prs)
    _section_tag(s, "Constats critiques")
    _slide_title(s, "Principaux risques identifies")

    # Priorite Major > Standard > Minor, complete jusqu'a 3 items meme si peu de Major
    top_risks = (majors_nc + standards_nc + minors_nc)[:3]
    y = 1.95
    rowH = 1.42
    if not top_risks:
        _add_text(s, "Aucun risque majeur identifie.", 0.6, y, 11.9, 0.5, size=14, italic=True, color=GRAY)
    for f in top_risks:
        _add_rect(s, 0.6, y, 11.9, rowH - 0.18, OFFWHITE, radius=0.03, line_color=LIGHTLINE, line_w=0.75)
        _add_ellipse(s, 0.85, y + 0.24, 0.62, 0.62, RED)
        _add_icon(s, "alert-triangle", 0.99, y + 0.38, 0.34, 0.34)
        _add_text(s, f.get("requirement_id", ""), 1.75, y + 0.14, 1.6, 0.9, size=10.5, bold=True, color=GRAY)
        _add_text(s, f.get("criterion", ""), 1.75, y + 0.34, 4.1, 0.7, size=14, bold=True)
        _add_text(s, f.get("justification", ""), 6.0, y + 0.14, 6.3, rowH - 0.4, size=11.5, color=GRAY,
                   anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.25)
        y += rowH

    _foot(s)

    # ------------------------------------------------- SLIDE 6+: RECOMMANDATIONS
    recommendations = report.get("recommendations", []) or []

    ROW_TOP = 1.95
    ROW_H = 0.92
    AVAILABLE_H_REC = SLIDE_H - ROW_TOP - 0.6
    PER_PAGE_REC = max(1, int(AVAILABLE_H_REC // ROW_H))

    rec_pages = [recommendations[i:i + PER_PAGE_REC] for i in range(0, len(recommendations), PER_PAGE_REC)] or [[]]

    for page_idx, chunk in enumerate(rec_pages):
        s = _blank_slide(prs)
        _section_tag(s, "Plan d'action")
        title = "Recommandations"
        if len(rec_pages) > 1:
            title += f"  ({page_idx + 1}/{len(rec_pages)})"
        _slide_title(s, title)

        if not recommendations:
            _add_text(s, "Aucune recommandation disponible.", 0.6, ROW_TOP, 11.9, 0.5, size=14, italic=True, color=GRAY)
        else:
            y = ROW_TOP
            base_num = page_idx * PER_PAGE_REC
            for i, rec in enumerate(chunk):
                _add_rect(s, 0.6, y, 11.9, ROW_H - 0.18, OFFWHITE, radius=0.1, line_color=LIGHTLINE, line_w=0.75)
                _add_ellipse(s, 0.85, y + (ROW_H - 0.18) / 2 - 0.22, 0.44, 0.44, NAVY)
                _add_text(s, str(base_num + i + 1), 0.85, y + (ROW_H - 0.18) / 2 - 0.22, 0.44, 0.44,
                           size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
                _add_text(s, rec, 1.55, y, 10.7, ROW_H - 0.18, size=13, color=NAVY_DARK,
                           anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.25)
                y += ROW_H

        _foot(s)


    Path(output_path).parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)
    return output_path


if __name__ == "__main__":
    import json, sys
    if len(sys.argv) < 3:
        print("Usage: python generate_report_pptx.py resultat.json sortie.pptx [nom_fournisseur]")
        sys.exit(1)
    with open(sys.argv[1], encoding="utf-8") as f:
        result = json.load(f)
    supplier = sys.argv[3] if len(sys.argv) > 3 else "[Nom du fournisseur]"
    out = generate_report_pptx(result, sys.argv[2], supplier_name=supplier)
    print("Genere:", out)
